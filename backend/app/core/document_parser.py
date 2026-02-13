"""
SmartRAG 智能文档解析引擎
支持: PDF, DOCX, TXT, Markdown, PPTX, XLSX, HTML
"""
import os
from typing import List, Optional
from dataclasses import dataclass, field
from loguru import logger


@dataclass
class ParsedDocument:
    """解析后的文档对象"""
    content: str
    pages: List[dict] = field(default_factory=list)
    # pages: [{"page": 1, "content": "...", "tables": [...]}]
    metadata: dict = field(default_factory=dict)
    tables: List[dict] = field(default_factory=list)


class DocumentParser:
    """多格式文档解析器"""

    SUPPORTED_TYPES = {
        ".pdf", ".docx", ".doc", ".txt", ".md",
        ".pptx", ".xlsx", ".csv", ".html", ".htm",
    }

    def parse(self, file_path: str) -> ParsedDocument:
        """主解析入口 - 根据文件类型路由到对应解析器"""
        ext = os.path.splitext(file_path)[1].lower()
        logger.info(f"Parsing document: {file_path} (type={ext})")

        if ext not in self.SUPPORTED_TYPES:
            raise ValueError(f"不支持的文件类型: {ext}")

        parser_map = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".doc": self._parse_docx,
            ".txt": self._parse_text,
            ".md": self._parse_markdown,
            ".pptx": self._parse_pptx,
            ".xlsx": self._parse_excel,
            ".csv": self._parse_csv,
            ".html": self._parse_html,
            ".htm": self._parse_html,
        }

        parser = parser_map.get(ext)
        if parser is None:
            raise ValueError(f"未实现的解析器: {ext}")

        result = parser(file_path)
        result.metadata["file_type"] = ext
        result.metadata["file_path"] = file_path
        result.metadata["file_name"] = os.path.basename(file_path)
        return result

    # ───────── PDF 解析 ─────────
    def _parse_pdf(self, file_path: str) -> ParsedDocument:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        pages = []
        full_text_parts = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                pages.append({
                    "page": i + 1,
                    "content": text,
                })
                full_text_parts.append(text)

        # 尝试提取表格(使用简单方式，生产环境推荐 camelot / tabula)
        tables = self._extract_pdf_tables(file_path)

        return ParsedDocument(
            content="\n\n".join(full_text_parts),
            pages=pages,
            tables=tables,
            metadata={
                "total_pages": len(reader.pages),
            },
        )

    def _extract_pdf_tables(self, file_path: str) -> List[dict]:
        """PDF 表格提取 (简化版)"""
        tables = []
        try:
            # 生产环境建议使用 camelot-py 或 pdfplumber
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_tables = page.extract_tables()
                    for t_idx, table in enumerate(page_tables):
                        if table:
                            # 将表格转为文本描述
                            header = table[0] if table else []
                            rows = table[1:] if len(table) > 1 else []
                            table_text = self._table_to_text(header, rows)
                            tables.append({
                                "page": i + 1,
                                "table_index": t_idx,
                                "content": table_text,
                                "raw": table,
                            })
        except ImportError:
            logger.warning("pdfplumber not installed, skipping table extraction")
        except Exception as e:
            logger.warning(f"Table extraction failed: {e}")
        return tables

    def _table_to_text(self, header: list, rows: list) -> str:
        """表格转自然语言描述"""
        lines = []
        header_clean = [str(h or "").strip() for h in header]
        lines.append("| " + " | ".join(header_clean) + " |")
        lines.append("| " + " | ".join(["---"] * len(header_clean)) + " |")
        for row in rows:
            row_clean = [str(c or "").strip() for c in row]
            lines.append("| " + " | ".join(row_clean) + " |")
        return "\n".join(lines)

    # ───────── DOCX 解析 ─────────
    def _parse_docx(self, file_path: str) -> ParsedDocument:
        from docx import Document as DocxDocument

        doc = DocxDocument(file_path)
        parts = []
        tables = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # 保留标题层级信息
                style = para.style.name if para.style else ""
                if "Heading" in style:
                    level = style.replace("Heading", "").strip() or "1"
                    parts.append(f"{'#' * int(level)} {text}")
                else:
                    parts.append(text)

        # 提取表格
        for i, table in enumerate(doc.tables):
            header = [cell.text.strip() for cell in table.rows[0].cells]
            rows = []
            for row in table.rows[1:]:
                rows.append([cell.text.strip() for cell in row.cells])
            table_text = self._table_to_text(header, rows)
            tables.append({
                "table_index": i,
                "content": table_text,
            })
            parts.append(f"\n[表格 {i + 1}]\n{table_text}\n")

        return ParsedDocument(
            content="\n\n".join(parts),
            tables=tables,
            metadata={"paragraph_count": len(doc.paragraphs)},
        )

    # ───────── TXT 解析 ─────────
    def _parse_text(self, file_path: str) -> ParsedDocument:
        import chardet
        with open(file_path, "rb") as f:
            raw = f.read()
        encoding = chardet.detect(raw)["encoding"] or "utf-8"
        content = raw.decode(encoding, errors="ignore")
        return ParsedDocument(content=content)

    # ───────── Markdown 解析 ─────────
    def _parse_markdown(self, file_path: str) -> ParsedDocument:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        return ParsedDocument(content=content)

    # ───────── PPTX 解析 ─────────
    def _parse_pptx(self, file_path: str) -> ParsedDocument:
        from pptx import Presentation

        prs = Presentation(file_path)
        pages = []
        all_text = []

        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
            slide_content = "\n".join(slide_texts)
            if slide_content:
                pages.append({"page": i + 1, "content": slide_content})
                all_text.append(f"[Slide {i + 1}]\n{slide_content}")

        return ParsedDocument(
            content="\n\n".join(all_text),
            pages=pages,
            metadata={"slide_count": len(prs.slides)},
        )

    # ───────── Excel 解析 ─────────
    def _parse_excel(self, file_path: str) -> ParsedDocument:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, data_only=True)
        all_text = []
        tables = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue

            header = [str(c or "") for c in rows[0]]
            data_rows = [[str(c or "") for c in r] for r in rows[1:]]
            table_text = self._table_to_text(header, data_rows)

            tables.append({
                "sheet": sheet_name,
                "content": table_text,
            })
            all_text.append(f"[Sheet: {sheet_name}]\n{table_text}")

        return ParsedDocument(
            content="\n\n".join(all_text),
            tables=tables,
            metadata={"sheet_count": len(wb.sheetnames)},
        )

    # ───────── CSV 解析 ─────────
    def _parse_csv(self, file_path: str) -> ParsedDocument:
        import csv
        import chardet

        with open(file_path, "rb") as f:
            raw = f.read()
        encoding = chardet.detect(raw)["encoding"] or "utf-8"

        lines = raw.decode(encoding).splitlines()
        reader = csv.reader(lines)
        rows = list(reader)

        if not rows:
            return ParsedDocument(content="")

        header = rows[0]
        data_rows = rows[1:]
        table_text = self._table_to_text(header, data_rows)

        return ParsedDocument(
            content=table_text,
            tables=[{"content": table_text}],
        )

    # ───────── HTML 解析 ─────────
    def _parse_html(self, file_path: str) -> ParsedDocument:
        from bs4 import BeautifulSoup
        from markdownify import markdownify

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html = f.read()

        soup = BeautifulSoup(html, "html.parser")

        # 移除无用标签
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        # HTML → Markdown (保留结构)
        md_content = markdownify(str(soup), heading_style="ATX")

        return ParsedDocument(
            content=md_content.strip(),
            metadata={"title": soup.title.string if soup.title else ""},
        )